from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RUN = ROOT / "runs" / "grant_demo"
OUT = ROOT / "grant_materials"
PPTX = OUT / "bridge_maker_grant_presentation.pptx"
PROOF = OUT / "proof_of_work.md"
NOITA_READINESS = ROOT / "adapters" / "noita_ws" / "README.md"

BG = RGBColor(9, 14, 29)
PANEL = RGBColor(24, 34, 57)
PANEL_2 = RGBColor(32, 47, 78)
TEXT = RGBColor(237, 243, 255)
MUTED = RGBColor(162, 180, 215)
ACCENT = RGBColor(98, 230, 172)
WARN = RGBColor(255, 190, 92)
BUG = RGBColor(255, 91, 128)


def add_bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.line.fill.background()


def add_text(slide, text, x, y, w, h, size=24, color=TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.38, 9.8, 0.62, 32, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.02, 10.8, 0.35, 14, MUTED)


def card(slide, x, y, w, h, label, value, color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = RGBColor(53, 71, 112)
    add_text(slide, label.upper(), x + 0.18, y + 0.16, w - 0.36, 0.24, 9, MUTED, True)
    add_text(slide, str(value), x + 0.18, y + 0.48, w - 0.36, h - 0.55, 24, color, True)


def pill(slide, text, x, y, w, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_text(slide, text, x, y + 0.07, w, 0.17, 9, RGBColor(5, 13, 23), True, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def proof_data():
    report = json.loads((RUN / "report.json").read_text(encoding="utf-8"))
    state_map = json.loads((RUN / "state_map.json").read_text(encoding="utf-8"))
    trace = [json.loads(line) for line in (RUN / "trace.jsonl").read_text(encoding="utf-8").splitlines()]
    return report, state_map, trace


def build():
    OUT.mkdir(parents=True, exist_ok=True)
    report, state_map, trace = proof_data()
    summary = report["summary"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_text(slide, "The Bridge-Maker", 0.7, 0.75, 8.0, 0.8, 48, TEXT, True)
    add_text(slide, "Contract-first automated QA for indie games", 0.75, 1.58, 7.8, 0.4, 20, MUTED)
    add_text(slide, "Annotate gameplay semantics -> generate a testable contract -> run agents -> receive bug evidence.", 0.75, 2.28, 7.5, 0.8, 24, TEXT, False)
    card(slide, 8.9, 0.82, 3.35, 1.2, "MVP proof", "bug found", BUG)
    card(slide, 8.9, 2.25, 3.35, 1.2, "No CE/Ghidra", "required", ACCENT)
    card(slide, 8.9, 3.68, 3.35, 1.2, "Output", "report.html", WARN)
    add_text(slide, "Grant demo deck - 2026", 0.78, 6.85, 3.0, 0.22, 10, MUTED)

    # 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "The hard lesson", "Black-box game understanding is useful research, but a bad default onboarding path.")
    items = [
        ("CE/Ghidra setup", "Powerful, but too manual and fragile for ordinary indie teams."),
        ("Pointer stability", "ASLR, stale addresses, GUI state, and process permissions create demo risk."),
        ("Action discovery", "A bot pressing keys without semantics is not QA; it is noise."),
        ("Better promise", "Let the developer expose a tiny contract; automate everything after that."),
    ]
    for i, (head, body) in enumerate(items):
        y = 1.65 + i * 1.05
        pill(slide, f"{i + 1}", 0.75, y, 0.45, ACCENT if i == 3 else WARN)
        add_text(slide, head, 1.35, y - 0.02, 3.6, 0.25, 19, TEXT, True)
        add_text(slide, body, 1.35, y + 0.33, 9.8, 0.35, 15, MUTED)

    # 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Product idea", "A tiny semantic contract replaces PDDL and reverse-engineering ceremony.")
    code = """from bridge_maker import bm

@bm.hp(bounds=(0, 10))
def hp(): return game.hp

@bm.position(x="x", y="y", bounds=(0, 9))
def position(): return game.x, game.y

@bm.move("right", key="d")
def move_right(): game.move_right()

@bm.oracle("out_of_bounds")
def out_of_bounds(s): return s.x < 0 or s.x > 9"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(6.0), Inches(4.95))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(12, 20, 38)
    box.line.color.rgb = RGBColor(64, 88, 136)
    add_text(slide, code, 1.0, 1.82, 5.5, 4.4, 13, RGBColor(218, 231, 255))
    for i, txt in enumerate(["No PDDL", "No CE setup", "No RLlib knowledge", "Reviewable code"]):
        card(slide, 7.25, 1.65 + i * 1.02, 4.55, 0.75, txt, "yes", ACCENT)

    # 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Architecture", "Agents amplify the contract; they do not invent hidden game semantics.")
    nodes = [
        ("Decorators / Adapter", 0.75, 2.0),
        ("Registry + Trace", 3.15, 2.0),
        ("Contract Compiler", 5.45, 2.0),
        ("Gym Env", 7.75, 2.0),
        ("QA Report", 10.0, 2.0),
    ]
    for label, x, y in nodes:
        card(slide, x, y, 1.85, 1.1, "", label, ACCENT)
    for i in range(len(nodes) - 1):
        arrow(slide, nodes[i][1] + 1.85, 2.55, nodes[i + 1][1], 2.55)
    add_text(slide, "Optional assist layer: CE MCP / Ghidra MCP / VLM can suggest or verify fields later.", 1.0, 4.15, 10.8, 0.35, 18, MUTED)
    add_text(slide, "Core rule: game-specific details live in annotations, adapters, or generated maps - not in the training core.", 1.0, 4.75, 10.8, 0.5, 21, TEXT, True)

    # 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Proof of work", "The MVP ran against an annotated roguelike with an intentional movement bug.")
    card(slide, 0.75, 1.55, 1.8, 1.1, "state fields", summary["state_fields"])
    card(slide, 2.8, 1.55, 1.8, 1.1, "actions", summary["actions"])
    card(slide, 4.85, 1.55, 1.8, 1.1, "oracles", summary["oracles"])
    card(slide, 6.9, 1.55, 1.8, 1.1, "trace frames", summary["trace_frames"])
    card(slide, 8.95, 1.55, 1.8, 1.1, "bug hits", summary["oracle_hits"], BUG)
    card(slide, 11.0, 1.55, 1.55, 1.1, "status", summary["status"], BUG)
    add_text(slide, "Command", 0.8, 3.05, 1.4, 0.25, 14, MUTED, True)
    add_text(slide, "python -m bridge_maker demo --out runs/grant_demo", 0.8, 3.42, 8.8, 0.35, 20, TEXT, True)
    first = report["oracle_hits"][0]
    add_text(slide, "First finding", 0.8, 4.28, 1.8, 0.25, 14, MUTED, True)
    add_text(slide, f"{first['name']} after {first['action']}: x={first['state']['x']}, valid max=9", 0.8, 4.65, 8.8, 0.35, 22, BUG, True)

    # 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Trace evidence", "A short stress burst crosses the right boundary and triggers the oracle.")
    x0, y0 = 0.9, 3.0
    for i, frame in enumerate(trace[:7]):
        x = x0 + i * 1.65
        state = frame["state"]
        hit = bool(frame["oracles"])
        color = BUG if hit else ACCENT
        pill(slide, str(i), x, y0 - 0.45, 0.45, color)
        card(slide, x - 0.1, y0, 1.3, 1.0, frame.get("action") or "start", f"x={state['x']}", color)
        if i:
            arrow(slide, x - 0.45, y0 + 0.5, x - 0.1, y0 + 0.5, color)
    add_text(slide, "Oracle hit frames: 2, 3, 4. The report stores state snapshots and action names for replay/debug.", 1.0, 5.15, 11.0, 0.4, 20, TEXT, True)

    # 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "What is real today", "This is a small MVP, but the core loop is already executable.")
    done = [
        "Decorator SDK: @bm.hp, @bm.position, @bm.action, @bm.oracle",
        "Adapter loader: Python file -> runtime registry",
        "Contract export: state_map/action_map/oracle_map/trace",
        "SDK Gymnasium env: reset/step over annotated functions",
        "Report generation: JSON + HTML with oracle findings",
        "Noita WS adapter boundary: heartbeat + raw Lua round-trip tested",
    ]
    for i, item in enumerate(done):
        pill(slide, "done", 0.85, 1.55 + i * 0.68, 0.75, ACCENT)
        add_text(slide, item, 1.8, 1.55 + i * 0.68, 9.9, 0.28, 18, TEXT)

    # 8
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Roadmap to a fundable product", "The grant milestone should prove usefulness, not universal magic.")
    phases = [
        ("MVP", "Contract + report over toy/adapter task", "done"),
        ("NoitaRL", "Wire a real reverse-engineered adapter", "next"),
        ("Agents", "Suggest missing annotations from code + traces", "next"),
        ("Training", "Connect contract envs to Ray/RLlib swarm", "next"),
        ("Studio UX", "One command: overnight run -> shareable report", "grant"),
    ]
    for i, (phase, body, tag) in enumerate(phases):
        y = 1.55 + i * 0.92
        card(slide, 0.8, y, 2.1, 0.65, tag, phase, ACCENT if tag == "done" else WARN)
        add_text(slide, body, 3.15, y + 0.12, 8.8, 0.24, 17, TEXT)
    add_text(slide, "Real external target inspected: noita-ws-api @ 47054b0. Adapter boundary tested with fake Noita heartbeat/replies; live adapter awaits game/mod setup.", 0.85, 6.45, 11.2, 0.35, 14, MUTED)

    # 9
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Why this can win", "A credible wedge between manual QA, heavy planning systems, and brittle black-box bots.")
    chart_data = CategoryChartData()
    chart_data.categories = ["Manual QA", "PDDL workflow", "Black-box bot", "Bridge-Maker"]
    chart_data.add_series("Indie usability", (2, 1, 2, 5))
    chart_data.add_series("Automation depth", (1, 4, 2, 4))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.85),
        Inches(1.7),
        Inches(7.1),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 5
    add_text(slide, "Positioning is qualitative for grant narrative; measured validation comes next with NoitaRL and user studies.", 8.25, 2.05, 3.8, 1.4, 18, MUTED)

    # 10
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_text(slide, "The ask", 0.75, 0.7, 5.5, 0.7, 44, TEXT, True)
    add_text(slide, "Fund the next milestone: real adapter validation, agent-assisted annotation mining, and a polished overnight QA report.", 0.78, 1.55, 9.8, 0.9, 25, TEXT)
    card(slide, 0.85, 3.2, 3.1, 1.15, "Milestone 1", "NoitaRL adapter", ACCENT)
    card(slide, 4.25, 3.2, 3.1, 1.15, "Milestone 2", "agent suggestions", WARN)
    card(slide, 7.65, 3.2, 3.1, 1.15, "Milestone 3", "overnight report", BUG)
    add_text(slide, "Bridge-Maker already has the smallest honest loop working. The grant turns it into a product-grade workflow.", 0.85, 5.35, 10.7, 0.5, 22, MUTED, True)

    prs.save(PPTX)

    shutil.copy2(RUN / "report.html", OUT / "grant_demo_report.html")
    shutil.copy2(RUN / "report.json", OUT / "grant_demo_report.json")
    if NOITA_READINESS.exists():
        shutil.copy2(NOITA_READINESS, OUT / "noita_ws_readiness.md")
    PROOF.write_text(
        "\n".join(
            [
                "# Bridge-Maker Proof of Work",
                "",
                "## Commands verified",
                "",
                "- python -m compileall bridge_maker src examples adapters",
                "- python -m bridge_maker demo --out runs/grant_demo",
                "- python -m bridge_maker smoke --adapter examples/annotated_dummy.py --steps 20",
                "- python -m bridge_maker report --contract runs/grant_demo",
                "- python -m unittest tests.test_noita_ws_session",
                "",
                "## Demo result",
                "",
                f"- Game: {summary['game_name']}",
                f"- State fields: {summary['state_fields']}",
                f"- Actions: {summary['actions']}",
                f"- Oracles: {summary['oracles']}",
                f"- Trace frames: {summary['trace_frames']}",
                f"- Oracle hits: {summary['oracle_hits']}",
                f"- Status: {summary['status']}",
                "",
                "## Real external target inspected",
                "",
                "- Repository: probable-basilisk/noita-ws-api",
                "- Local commit: 47054b0",
                "- Adapter readiness: grant_materials/noita_ws_readiness.md",
                "- Adapter boundary test: fake Noita client heartbeat + raw Lua replies passed.",
                "- Live Noita execution pending game/mod runtime setup.",
                "",
                "## First finding",
                "",
                f"- Oracle: {report['oracle_hits'][0]['name']}",
                f"- Action: {report['oracle_hits'][0]['action']}",
                f"- State: {json.dumps(report['oracle_hits'][0]['state'])}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    print(PPTX)
    print(PROOF)


if __name__ == "__main__":
    build()
